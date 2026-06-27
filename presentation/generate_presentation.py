#!/usr/bin/env python3
"""Generate Cultural & Cross-Cultural Awareness sociology presentation."""

from __future__ import annotations

import io
from pathlib import Path

import requests
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).parent / "Cultural_Cross_Cultural_Awareness_Presentation.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BLACK = RGBColor(0x0A, 0x0A, 0x0A)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
MID = RGBColor(0x2E, 0x2E, 0x2E)
GRAY = RGBColor(0x88, 0x88, 0x88)
LIGHT = RGBColor(0xCC, 0xCC, 0xCC)
WHITE = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT = RGBColor(0xE8, 0xE8, 0xE8)

ASSETS = Path(__file__).parent / "assets"


def download_asset(name: str, url: str) -> Path | None:
    path = ASSETS / name
    if path.exists() and path.stat().st_size > 500:
        return path
    try:
        resp = requests.get(url, timeout=20, headers={"User-Agent": "Mozilla/5.0"})
        resp.raise_for_status()
        path.write_bytes(resp.content)
        return path
    except Exception:
        return None


def make_gradient_bg(width: int = 1920, height: int = 1080) -> io.BytesIO:
    img = Image.new("RGB", (width, height))
    draw = ImageDraw.Draw(img)
    for y in range(height):
        t = y / height
        v = int(18 + (55 - 18) * t)
        draw.line([(0, y), (width, y)], fill=(v, v, v))
    for i in range(-height, width, 48):
        draw.line([(i, 0), (i + height, height)], fill=(30, 30, 30), width=1)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def set_slide_bg(slide, prs: Presentation) -> None:
    buf = make_gradient_bg()
    slide.shapes.add_picture(buf, 0, 0, width=SLIDE_W, height=SLIDE_H)


def add_accent_bar(slide, y=Inches(0.55), width=Inches(2.2)) -> None:
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.85), y, width, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = WHITE
    bar.line.fill.background()


def add_footer(slide, text: str = "Educational Dialogues on Cultural Awareness") -> None:
    box = slide.shapes.add_textbox(Inches(0.85), Inches(7.05), Inches(8), Inches(0.35))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY
    p.font.name = "Arial"


def style_title(tf, text: str, size: int = 40, bold: bool = True) -> None:
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text.upper()
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = WHITE
    p.font.name = "Arial Black"
    p.alignment = PP_ALIGN.LEFT


def add_label_box(slide, left, top, width, height, label: str, body_lines: list[str], label_size=11, body_size=14):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = MID
    shape.line.color.rgb = LIGHT
    shape.line.width = Pt(0.75)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(10)
    tf.margin_bottom = Pt(10)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.text = label.upper()
    p.font.size = Pt(label_size)
    p.font.bold = True
    p.font.color.rgb = LIGHT
    p.font.name = "Arial"
    p.space_after = Pt(6)

    for line in body_lines:
        bp = tf.add_paragraph()
        bp.text = f"• {line}" if not line.startswith("•") else line
        bp.font.size = Pt(body_size)
        bp.font.color.rgb = WHITE
        bp.font.name = "Arial"
        bp.space_after = Pt(4)


def add_takeaway(slide, text: str) -> None:
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(6.35), Inches(7.8), Inches(0.55)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = BLACK
    box.line.color.rgb = WHITE
    box.line.width = Pt(1)
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = f"KEY TAKEAWAY: {text}"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER


def add_image_or_placeholder(slide, path: Path | None, left, top, width, height, caption: str = "") -> None:
    if path and path.exists():
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    else:
        ph = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
        ph.fill.solid()
        ph.fill.fore_color.rgb = DARK
        ph.line.color.rgb = GRAY
        ph.line.width = Pt(1)
        tf = ph.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = caption or "Visual"
        p.font.size = Pt(10)
        p.font.color.rgb = GRAY
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER


def add_icon_circle(slide, left, top, size, symbol: str) -> None:
    circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = BLACK
    circle.line.color.rgb = WHITE
    circle.line.width = Pt(1.5)
    tf = circle.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = symbol
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def concept_slide(
    prs: Presentation,
    title: str,
    definition: list[str],
    example: list[str],
    takeaway: str,
    visuals: list[tuple],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, prs)
    add_accent_bar(slide)
    add_footer(slide)

    title_box = slide.shapes.add_textbox(Inches(0.85), Inches(0.75), Inches(8), Inches(0.9))
    style_title(title_box.text_frame, title, size=36)

    add_label_box(slide, Inches(0.85), Inches(1.85), Inches(5.9), Inches(1.55), "Definition", definition)
    add_label_box(slide, Inches(0.85), Inches(3.55), Inches(5.9), Inches(2.35), "Example", example)
    add_takeaway(slide, takeaway)

    for path, caption, l, t, w, h in visuals:
        add_image_or_placeholder(slide, path, Inches(l), Inches(t), Inches(w), Inches(h), caption)


def local_asset(name: str) -> Path | None:
    path = ASSETS / name
    return path if path.exists() else None


def fetch_assets() -> dict[str, Path | None]:
    names = [
        "world_map.png", "us_flag.png", "german_flag.png", "cherokee_pattern.png",
        "migration_map.png", "diverse_people.png", "sushi.png",
    ]
    return {n: local_asset(n) for n in names}


def build_presentation() -> Path:
    ASSETS.mkdir(parents=True, exist_ok=True)
    assets = fetch_assets()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, prs)
    add_accent_bar(slide, y=Inches(2.35), width=Inches(4.5))

    title_box = slide.shapes.add_textbox(Inches(0.85), Inches(1.35), Inches(11.5), Inches(1.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "EDUCATIONAL DIALOGUES ON"
    p.font.size = Pt(22)
    p.font.color.rgb = LIGHT
    p.font.name = "Arial"
    p2 = tf.add_paragraph()
    p2.text = "CULTURAL AND CROSS-CULTURAL AWARENESS"
    p2.font.size = Pt(38)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.font.name = "Arial Black"

    meta = slide.shapes.add_textbox(Inches(0.85), Inches(3.55), Inches(6), Inches(1.6))
    mtf = meta.text_frame
    for i, line in enumerate(["Palak Patel", "June 26, 2026", "United States 🇺🇸"]):
        mp = mtf.paragraphs[0] if i == 0 else mtf.add_paragraph()
        mp.text = line
        mp.font.size = Pt(20 if i == 0 else 16)
        mp.font.bold = i == 0
        mp.font.color.rgb = WHITE if i == 0 else LIGHT
        mp.font.name = "Arial"
        mp.space_after = Pt(8)

    add_image_or_placeholder(slide, assets.get("us_flag.png"), Inches(9.8), Inches(1.2), Inches(2.4), Inches(1.5), "U.S. Flag")
    add_image_or_placeholder(slide, assets.get("world_map.png"), Inches(7.2), Inches(3.2), Inches(5.5), Inches(3.5), "World Map")

    # Introduction
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, prs)
    add_accent_bar(slide)
    add_footer(slide)
    title_box = slide.shapes.add_textbox(Inches(0.85), Inches(0.75), Inches(8), Inches(0.8))
    style_title(title_box.text_frame, "Introduction", size=34)

    add_label_box(slide, Inches(0.85), Inches(1.75), Inches(5.5), Inches(1.4), "What is culture?", [
        "Shared beliefs, values, behaviors, and symbols",
        "Passed between generations and groups",
        "Shapes identity and everyday life",
    ])
    add_label_box(slide, Inches(0.85), Inches(3.35), Inches(5.5), Inches(1.4), "Why it matters", [
        "Builds empathy across differences",
        "Reduces bias and misunderstanding",
        "Prepares us for a connected world",
    ])

    quote = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(5.0), Inches(7.2), Inches(1.2))
    quote.fill.solid()
    quote.fill.fore_color.rgb = BLACK
    quote.line.color.rgb = WHITE
    quote.line.width = Pt(1)
    qtf = quote.text_frame
    qtf.vertical_anchor = MSO_ANCHOR.MIDDLE
    qp = qtf.paragraphs[0]
    qp.text = '"Culture is the widening of the mind and of the spirit." — Jawaharlal Nehru'
    qp.font.size = Pt(14)
    qp.font.italic = True
    qp.font.color.rgb = ACCENT
    qp.font.name = "Arial"
    qp.alignment = PP_ALIGN.CENTER

    add_image_or_placeholder(slide, assets.get("world_map.png"), Inches(7.5), Inches(1.75), Inches(2.2), Inches(2.2), "🌍")
    add_icon_circle(slide, Inches(10.2), Inches(2.0), Inches(1.1), "🤝")
    add_icon_circle(slide, Inches(10.2), Inches(3.5), Inches(1.1), "💬")
    add_icon_circle(slide, Inches(10.2), Inches(5.0), Inches(1.1), "🌐")

    concepts = [
        {"title": "Culture", "definition": ["Shared values, beliefs, customs, and behaviors", "Shape a group's way of life and identity"], "example": ["Krystal's identity blends German and Cherokee heritage", "She moved from Los Angeles to the Southeastern U.S.", "Multiple cultural influences shape how she sees the world"], "takeaway": "Culture is learned, shared, and always evolving.", "visuals": [(assets.get("german_flag.png"), "German Flag", 7.2, 1.75, 1.6, 1.0), (assets.get("cherokee_pattern.png"), "Cherokee Heritage", 9.0, 1.75, 1.8, 1.0), (assets.get("migration_map.png"), "CA to Southeast U.S.", 7.2, 3.0, 3.6, 1.8), (assets.get("diverse_people.png"), "Diverse Community", 7.2, 5.0, 3.6, 1.1)]},
        {"title": "Cultural Universal", "definition": ["A trait or practice found in every known society", "Form varies, but the institution appears everywhere"], "example": ["Family, language, marriage, religion, and music", "Every society creates ways to connect and belong"], "takeaway": "Humans share core institutions — expression differs.", "visuals": [(None, "Family", 7.2, 1.75, 1.5, 1.5), (None, "Language", 9.0, 1.75, 1.5, 1.5), (None, "Music", 7.2, 3.5, 1.5, 1.5), (None, "Marriage & Religion", 9.0, 3.5, 1.8, 1.5)]},
        {"title": "Ethnocentrism", "definition": ["Judging another culture by your own standards", "Often treats your culture as superior or 'normal'"], "example": ["Calling insect-based diets 'gross'", "Because they differ from your culture's food norms"], "takeaway": "My culture is the center — not the universal standard.", "visuals": [(None, "⚖️ My Way vs. Their Way", 7.2, 1.75, 3.6, 2.5), (None, "🍽️ Food Norms", 7.2, 4.5, 3.6, 1.6)]},
        {"title": "Cultural Relativism", "definition": ["Evaluate a culture on its own values and context", "Suspend judgment based on outside standards"], "example": ["Understanding arranged marriages within their social context", "Asking why practices exist before labeling them"], "takeaway": "Understand a culture by its own standards.", "visuals": [(None, "🔍 Context First", 7.2, 1.75, 3.6, 2.2), (None, "🌍 Local Meaning", 7.2, 4.2, 3.6, 2.0)]},
        {"title": "Language", "definition": ["A system of symbols used to communicate culture", "Shapes thought, identity, and social interaction"], "example": ["English, Spanish, and American Sign Language", "Cherokee language preservation efforts"], "takeaway": "Language carries culture — and can be lost or revived.", "visuals": [(None, "Speech", 7.2, 1.75, 1.8, 1.8), (None, "A · B · C", 9.3, 1.75, 1.5, 1.8), (None, "Cherokee Syllabary", 7.2, 3.8, 3.6, 1.8)]},
        {"title": "Non-verbal Communication", "definition": ["Communication without words", "Gestures, expressions, tone, eye contact, space"], "example": ["A thumbs-up, nodding, or smiling", "Personal space norms differ across cultures"], "takeaway": "Messages travel through body language too.", "visuals": [(None, "👍 Thumbs Up", 7.2, 1.75, 1.6, 1.6), (None, "😊 Expression", 9.1, 1.75, 1.6, 1.6), (None, "↔️ Personal Space", 7.2, 3.6, 3.5, 2.0)]},
        {"title": "Norms", "definition": ["Shared rules and expectations that guide behavior", "Help societies run smoothly and predictably"], "example": ["Waiting in line", "Saying 'thank you'"], "takeaway": "Norms = rules for everyday behavior.", "visuals": [(None, "Courtesy", 7.2, 1.75, 2.0, 2.0), (None, "🚶 Queue Culture", 9.5, 1.75, 1.5, 2.0), (None, "✓ Expected Conduct", 7.2, 4.0, 3.8, 1.8)]},
        {"title": "Types of Norms", "definition": ["Folkways: everyday customs", "Mores: moral expectations · Laws: formal rules · Taboos: strong prohibitions"], "example": ["Folkway: shaking hands", "More: not cheating · Law: speed limits · Taboo: incest"], "takeaway": "Not all norms carry the same weight.", "visuals": [(None, "🤝 Folkway", 7.2, 1.75, 1.7, 1.4), (None, "⚖️ More", 9.1, 1.75, 1.7, 1.4), (None, "📜 Law", 7.2, 3.3, 1.7, 1.4), (None, "🚫 Taboo", 9.1, 3.3, 1.7, 1.4)]},
        {"title": "Cultural Values", "definition": ["Shared ideas about what is good and desirable", "Guide priorities, goals, and social choices"], "example": ["Independence in the United States", "Collectivism in many Asian cultures"], "takeaway": "Values = ideas about what matters most.", "visuals": [(None, "🇺🇸 Individualism", 7.2, 1.75, 1.7, 2.0), (None, "🌏 Collectivism", 9.1, 1.75, 1.7, 2.0), (None, "⭐ Shared Priorities", 7.2, 4.0, 3.6, 1.8)]},
        {"title": "Sanctions", "definition": ["Rewards or punishments that enforce norms", "Can be positive/negative and formal/informal"], "example": ["Praise for volunteering (positive informal)", "Fine for speeding (negative formal)"], "takeaway": "Sanctions reward conformity and discourage violation.", "visuals": [(None, "🏆 Positive", 7.2, 1.75, 1.7, 1.6), (None, "⚠️ Negative", 9.1, 1.75, 1.7, 1.6), (None, "Formal vs Informal", 7.2, 3.6, 3.6, 2.0)]},
        {"title": "Culture War", "definition": ["Conflict over competing values, beliefs, or social issues", "Groups fight to define what society should accept"], "example": ["Debates over abortion, same-sex marriage, school curricula"], "takeaway": "Deep value clashes can divide communities.", "visuals": [(None, "⚔️ Competing Values", 7.2, 1.75, 3.6, 2.2), (None, "📢 Public Debate", 7.2, 4.2, 3.6, 2.0)]},
        {"title": "Dominant Ideology", "definition": ["Beliefs promoted by the most powerful groups", "Often accepted as 'normal' or common sense"], "example": ["The 'American Dream': hard work leads to success"], "takeaway": "Power shapes which ideas feel 'natural.'", "visuals": [(None, "🏛️ Power & Ideas", 7.2, 1.75, 3.6, 2.0), (None, "📈 Merit Narrative", 7.2, 4.0, 3.6, 2.0)]},
        {"title": "Cultural Variation", "definition": ["Differences in beliefs, customs, and behaviors", "Appear among societies and within them"], "example": ["Eating with chopsticks vs. forks", "Regional accents, dress, and traditions"], "takeaway": "Variation is normal — not a defect.", "visuals": [(None, "🥢 Chopsticks", 7.2, 1.75, 1.7, 1.8), (None, "🍴 Fork", 9.1, 1.75, 1.7, 1.8), (None, "🗺️ Global Differences", 7.2, 3.8, 3.6, 1.8)]},
        {"title": "Subculture", "definition": ["Distinct beliefs/behaviors within a larger culture", "Different — but still compatible with mainstream society"], "example": ["Gamers, skateboarders, medical students"], "takeaway": "Sub = under the umbrella of the larger culture.", "visuals": [(None, "🎮 Gamers", 7.2, 1.75, 1.7, 1.5), (None, "🛹 Skaters", 9.1, 1.75, 1.7, 1.5), (None, "☂️ Mainstream Umbrella", 7.2, 3.5, 3.6, 2.2)]},
        {"title": "Counterculture", "definition": ["Values and norms that challenge dominant culture", "Actively reject or oppose mainstream ideals"], "example": ["1960s hippie movement rejecting consumerism"], "takeaway": "Counter = against the dominant culture.", "visuals": [(None, "✌️ 1960s Movement", 7.2, 1.75, 3.6, 2.2), (None, "🔄 Reject Mainstream", 7.2, 4.2, 3.6, 2.0)]},
        {"title": "Culture Shock", "definition": ["Confusion, anxiety, or disorientation in a new culture", "Common when customs and cues feel unfamiliar"], "example": ["Krystal adjusting from Los Angeles to the Southern U.S.", "New accents, pace, food, and social expectations"], "takeaway": "Adaptation takes time — discomfort is part of learning.", "visuals": [(assets.get("migration_map.png"), "Los Angeles to Southeast", 7.2, 1.75, 3.6, 2.0), (None, "✈️ Transition", 7.2, 4.0, 3.6, 2.0)]},
        {"title": "Innovation", "definition": ["Creation of new ideas, technologies, or practices", "Can reshape how societies live and communicate"], "example": ["Smartphones and social media platforms"], "takeaway": "Innovation = creating something new.", "visuals": [(None, "📱 Smartphones", 7.2, 1.75, 1.8, 2.0), (None, "💡 New Ideas", 9.3, 1.75, 1.5, 2.0), (None, "🚀 Cultural Change", 7.2, 4.0, 3.6, 1.8)]},
        {"title": "Diffusion", "definition": ["Spread of cultural traits across societies", "Ideas, foods, and styles travel through contact"], "example": ["Sushi becoming popular worldwide"], "takeaway": "Diffusion = spreading what already exists.", "visuals": [(assets.get("sushi.png"), "Sushi Worldwide", 7.2, 1.75, 3.6, 2.0), (assets.get("world_map.png"), "Spreading Traits", 7.2, 4.0, 3.6, 2.0)]},
    ]

    for c in concepts:
        concept_slide(prs, c["title"], c["definition"], c["example"], c["takeaway"], c["visuals"])

    # Material vs Non-material
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, prs)
    add_accent_bar(slide)
    add_footer(slide)
    title_box = slide.shapes.add_textbox(Inches(0.85), Inches(0.75), Inches(10), Inches(0.8))
    style_title(title_box.text_frame, "Material vs. Non-material Culture", size=30)
    add_label_box(slide, Inches(0.85), Inches(1.85), Inches(5.5), Inches(2.0), "Material Culture", ["Physical objects and technology created by society", "Buildings, clothing, cars, Cherokee stone tools"])
    add_label_box(slide, Inches(0.85), Inches(4.05), Inches(5.5), Inches(2.0), "Non-material Culture", ["Intangible beliefs, values, norms, and traditions", "Religion, Pentecostal beliefs, laws, morals"])
    add_takeaway(slide, "Objects vs. ideas — both carry culture.")
    add_image_or_placeholder(slide, None, Inches(7.0), Inches(1.85), Inches(1.8), Inches(1.8), "🪓 Stone Axe")
    add_image_or_placeholder(slide, None, Inches(9.0), Inches(1.85), Inches(1.8), Inches(1.8), "🏠 Buildings")
    add_image_or_placeholder(slide, None, Inches(7.0), Inches(4.05), Inches(1.8), Inches(1.8), "⛪ Beliefs")
    add_image_or_placeholder(slide, None, Inches(9.0), Inches(4.05), Inches(1.8), Inches(1.8), "📖 Traditions")

    concept_slide(prs, "Cultural Lag", ["Non-material culture adapts slower than material culture", "Laws and ethics trail behind new technology"], ["AI and gene editing outpacing regulation", "Society debates ethics while tech advances quickly"], "Society's values and laws lag behind technology.", [(None, "🤖 AI Ethics", 7.2, 1.75, 1.8, 1.8), (None, "🧬 Gene Editing", 9.2, 1.75, 1.8, 1.8), (None, "⏳ The Gap", 7.2, 3.8, 3.8, 2.0)])

    # High-yield distinctions
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, prs)
    add_accent_bar(slide)
    add_footer(slide)
    title_box = slide.shapes.add_textbox(Inches(0.85), Inches(0.75), Inches(10), Inches(0.8))
    style_title(title_box.text_frame, "High-Yield Distinctions", size=30)
    pairs = [
        ("Ethnocentrism", "My culture is the standard."),
        ("Cultural Relativism", "Understand by its own standards."),
        ("Subculture", "Different but compatible."),
        ("Counterculture", "Challenges mainstream society."),
        ("Innovation", "Creating something new."),
        ("Diffusion", "Spreading what exists."),
        ("Norms", "Rules for behavior."),
        ("Values", "Ideas about importance."),
        ("Sanctions", "Rewards and punishments."),
    ]
    y = 1.65
    for left_label, right_text in pairs:
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(y), Inches(11.0), Inches(0.48))
        box.fill.solid()
        box.fill.fore_color.rgb = DARK
        box.line.color.rgb = GRAY
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = f"{left_label}  →  {right_text}"
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.font.name = "Arial"
        y += 0.55

    # Summary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, prs)
    add_accent_bar(slide)
    add_footer(slide)
    title_box = slide.shapes.add_textbox(Inches(0.85), Inches(0.75), Inches(8), Inches(0.8))
    style_title(title_box.text_frame, "Summary", size=34)
    add_label_box(slide, Inches(0.85), Inches(1.75), Inches(11.0), Inches(2.2), "Major Concepts", [
        "Culture shapes identity through shared symbols, norms, and values",
        "Awareness requires cultural relativism over ethnocentrism",
        "Subcultures, countercultures, and variation show culture's complexity",
        "Innovation, diffusion, and cultural lag explain how culture changes",
    ])
    add_label_box(slide, Inches(0.85), Inches(4.2), Inches(11.0), Inches(1.5), "Final Takeaways", [
        "Listen before judging — context matters",
        "Recognize your own cultural lens",
        "Cross-cultural awareness builds empathy and respect",
    ])

    # References
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, prs)
    add_accent_bar(slide)
    title_box = slide.shapes.add_textbox(Inches(0.85), Inches(0.75), Inches(8), Inches(0.8))
    style_title(title_box.text_frame, "References", size=34)
    add_label_box(slide, Inches(0.85), Inches(1.75), Inches(11.0), Inches(4.5), "Sources", [
        "Course material — Sociology: Culture and Society unit",
        "Wikimedia Commons — flags, icons, and map imagery (public domain / CC)",
        "Educational Dialogues on Cultural and Cross-Cultural Awareness",
        "Prepared by Palak Patel — June 26, 2026",
    ], body_size=13)

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build_presentation()
    from pptx import Presentation as P
    print(f"Created: {path}")
    print(f"Slides: {len(P(str(path)).slides)}")
